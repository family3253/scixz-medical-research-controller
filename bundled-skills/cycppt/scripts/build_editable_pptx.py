#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import math
from pathlib import Path

from PIL import Image, ImageFont, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

IMG_W = 2560
IMG_H = 1440
SLIDE_W_IN = 16
SLIDE_H_IN = 9
PX_PER_IN = IMG_W / SLIDE_W_IN
PT_PER_PX = 72 / PX_PER_IN
CJK_FONT = "Arial Unicode MS"
LATIN_FONT = "Arial"
FONT_PATHS = {
    CJK_FONT: [
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/NISC18030.ttf",
    ],
    LATIN_FONT: [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
    ],
}


def is_chinese(text: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def bbox_for(item: dict) -> tuple[float, float, float, float]:
    x1, y1, x2, y2 = item["bbox"]
    return float(x1), float(y1), float(x2), float(y2)


def poly_height_for(item: dict) -> float:
    poly = item.get("poly") or []
    if len(poly) >= 4:
        try:
            pts = [(float(x), float(y)) for x, y in poly[:4]]
            left = math.dist(pts[0], pts[3])
            right = math.dist(pts[1], pts[2])
            height = (left + right) / 2
            if height > 1:
                return height
        except (TypeError, ValueError):
            pass
    x1, y1, x2, y2 = bbox_for(item)
    return max(y2 - y1, 1.0)


def role_for(item: dict) -> str:
    text = item["text"].strip()
    x1, y1, x2, y2 = bbox_for(item)
    w = x2 - x1
    h = y2 - y1
    if x1 > 2320 and y1 < 150:
        return "logo"
    if y1 < 150:
        return "header"
    if y1 > 1260:
        return "footer"
    if y1 < 450 and h > 48 and w > 120:
        return "title"
    if text in {"1", "2", "3", "4", "5"} and w < 70 and h > 28:
        return "badge"
    if "Take-home" in text or "Scientific question" in text or "Method take-home" in text:
        return "takehome"
    if h > 36 and w > 100:
        return "heading"
    return "body"


def font_family_for(text: str) -> str:
    return CJK_FONT if is_chinese(text) else LATIN_FONT


def font_file_for(family: str) -> str | None:
    for raw in FONT_PATHS.get(family, []):
        path = Path(raw)
        if path.exists():
            return str(path)
    return None


def measure_text_px(text: str, font_pt: float, family: str) -> float:
    font_file = font_file_for(family)
    if font_file:
        font_px = max(1, int(round(font_pt / PT_PER_PX)))
        font = ImageFont.truetype(font_file, font_px)
        return float(font.getlength(text))
    weight = 0.0
    for ch in text:
        if "\u4e00" <= ch <= "\u9fff":
            weight += 0.92
        elif ch.isspace():
            weight += 0.32
        elif ch.isascii() and ch.isalnum():
            weight += 0.58
        else:
            weight += 0.45
    return weight * font_pt / PT_PER_PX


def max_font_for(item: dict, role: str) -> float:
    if role == "header":
        return 21.0
    if role == "logo":
        return 18.0
    if role == "footer":
        return 15.0
    if role == "title":
        return 48.0
    if role == "takehome":
        return 30.0
    if role in {"heading", "badge"}:
        return 24.0
    return 18.0


def min_font_for(role: str) -> float:
    if role in {"footer", "logo"}:
        return 5.0
    if role == "body":
        return 6.0
    return 7.0


def font_size_for(item: dict, box_width_px: float) -> float:
    text = item["text"].strip()
    role = role_for(item)
    family = font_family_for(text)
    height_pt = poly_height_for(item) * PT_PER_PX / 0.78
    target = min(height_pt, max_font_for(item, role))
    if role == "badge":
        target = min(target, 18.0)

    limit = max(box_width_px * 0.96, 1.0)
    if measure_text_px(text, target, family) > limit:
        lo = min_font_for(role)
        hi = target
        for _ in range(12):
            mid = (lo + hi) / 2
            if measure_text_px(text, mid, family) <= limit:
                lo = mid
            else:
                hi = mid
        target = lo
    return round(max(min_font_for(role), target), 1)


def expand_box(item: dict) -> tuple[float, float, float, float]:
    x1, y1, x2, y2 = bbox_for(item)
    w = max(x2 - x1, 1.0)
    h = max(y2 - y1, 1.0)
    line_h = poly_height_for(item)
    pad_x = max(8.0, min(24.0, line_h * 0.35))
    box_w_px = w + pad_x
    font_pt = font_size_for(item, box_w_px)
    font_px = font_pt / PT_PER_PX
    box_h_px = max(h + max(8.0, line_h * 0.45), font_px * 1.28)
    left_px = max(x1 - pad_x / 2, 0)
    top_px = max(y1 - (box_h_px - h) / 2, 0)
    left = left_px / IMG_W * SLIDE_W_IN
    top = top_px / IMG_H * SLIDE_H_IN
    width = max(box_w_px / IMG_W * SLIDE_W_IN, 0.08)
    height = max(box_h_px / IMG_H * SLIDE_H_IN, 0.08)
    return left, top, width, height


def fallback_color_for(item: dict) -> RGBColor:
    text = item["text"].strip()
    role = role_for(item)
    if role == "takehome":
        return RGBColor(150, 35, 35)
    if role in {"header", "title", "heading"}:
        return RGBColor(15, 35, 70)
    if role == "footer":
        return RGBColor(90, 90, 90)
    if text.startswith("Problem") or text in {"1", "2", "3", "4"}:
        return RGBColor(255, 255, 255)
    return RGBColor(20, 35, 60)


def median_rgb(pixels: list[tuple[int, int, int]]) -> tuple[int, int, int] | None:
    if not pixels:
        return None
    channels = list(zip(*pixels))
    return tuple(int(sorted(ch)[len(ch) // 2]) for ch in channels)


def crop_for_bbox(image: Image.Image, item: dict, pad: int = 1) -> Image.Image:
    x1, y1, x2, y2 = bbox_for(item)
    box = (
        max(int(math.floor(x1)) - pad, 0),
        max(int(math.floor(y1)) - pad, 0),
        min(int(math.ceil(x2)) + pad, image.width),
        min(int(math.ceil(y2)) + pad, image.height),
    )
    return image.crop(box).convert("RGB")


def sampled_text_color(item: dict, source_image: Image.Image | None, clean_image: Image.Image | None) -> RGBColor | None:
    if source_image is None:
        return None
    source_crop = crop_for_bbox(source_image, item)
    text_pixels: list[tuple[int, int, int]] = []
    if clean_image is not None:
        clean_crop = crop_for_bbox(clean_image, item).resize(source_crop.size)
        for src, clean in zip(source_crop.getdata(), clean_crop.getdata()):
            diff = sum(abs(src[i] - clean[i]) for i in range(3))
            if diff > 36:
                text_pixels.append(src)
    if len(text_pixels) < 8:
        stat = ImageStat.Stat(source_crop)
        bg = tuple(int(v) for v in stat.median[:3])
        for src in source_crop.getdata():
            diff = sum(abs(src[i] - bg[i]) for i in range(3))
            if diff > 70:
                text_pixels.append(src)
    rgb = median_rgb(text_pixels)
    if rgb is None:
        return None
    return RGBColor(*rgb)


def color_for(item: dict, source_image: Image.Image | None, clean_image: Image.Image | None) -> RGBColor:
    sampled = sampled_text_color(item, source_image, clean_image)
    if sampled is not None:
        return sampled
    return fallback_color_for(item)


def merge_text(left: str, right: str, gap: float) -> str:
    if gap > 70:
        return f"{left} | {right}"
    if left and right and (left[-1].isascii() or right[0].isascii()):
        return f"{left} {right}"
    return f"{left}{right}"


def can_merge_text_items(left: dict, right: dict) -> bool:
    left_text = left.get("text", "").strip()
    right_text = right.get("text", "").strip()
    if not left_text or not right_text:
        return False
    if role_for(left) not in {"header", "title"} or role_for(right) not in {"header", "title"}:
        return False
    lx1, ly1, lx2, ly2 = bbox_for(left)
    rx1, ry1, rx2, ry2 = bbox_for(right)
    lh = ly2 - ly1
    rh = ry2 - ry1
    centers_close = abs(((ly1 + ly2) / 2) - ((ry1 + ry2) / 2)) <= max(18, min(lh, rh) * 0.45)
    gap = rx1 - lx2
    return centers_close and -60 <= gap <= 110


def merge_adjacent_heading_items(items: list[dict]) -> list[dict]:
    ordered = sorted(items, key=lambda it: (bbox_for(it)[1], bbox_for(it)[0]))
    merged: list[dict] = []
    for item in ordered:
        if merged and can_merge_text_items(merged[-1], item):
            prev = merged[-1]
            px1, py1, px2, py2 = bbox_for(prev)
            x1, y1, x2, y2 = bbox_for(item)
            prev["text"] = merge_text(prev["text"].strip(), item["text"].strip(), x1 - px2)
            prev["bbox"] = [min(px1, x1), min(py1, y1), max(px2, x2), max(py2, y2)]
            prev["poly"] = [[prev["bbox"][0], prev["bbox"][1]], [prev["bbox"][2], prev["bbox"][1]], [prev["bbox"][2], prev["bbox"][3]], [prev["bbox"][0], prev["bbox"][3]]]
            prev["score"] = min(float(prev.get("score", 1.0)), float(item.get("score", 1.0)))
            prev["_merged"] = True
            continue
        merged.append(dict(item))
    return sorted(merged, key=lambda it: (bbox_for(it)[1], bbox_for(it)[0]))


def add_textbox(slide, item: dict, source_image: Image.Image | None = None, clean_image: Image.Image | None = None) -> None:
    text = item["text"].strip()
    if not text:
        return
    left, top, width, height = expand_box(item)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font_family_for(text)
    run.font.size = Pt(font_size_for(item, (width / SLIDE_W_IN) * IMG_W))
    run.font.color.rgb = color_for(item, source_image, clean_image)
    role = role_for(item)
    if role in {"header", "title", "heading", "takehome"}:
        run.font.bold = True


def slide_ids_from_manifest(manifest_path: Path | None, ocr_dir: Path | None = None) -> list[str]:
    if manifest_path:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        ids = [slide["slide_id"] for slide in manifest.get("slides", [])]
        if ids:
            return ids
    if ocr_dir:
        return [p.name.replace("_ocr.json", "") for p in sorted(ocr_dir.glob("slide*_ocr.json"))]
    return []


def load_notes(notes_path: Path | None) -> dict[str, str]:
    if not notes_path:
        return {}
    data = json.loads(notes_path.read_text(encoding="utf-8"))
    if isinstance(data, dict) and "slides" not in data:
        return {str(k): str(v) for k, v in data.items()}
    notes = {}
    for slide in data.get("slides", []):
        text = slide.get("speaker_notes_zh") or slide.get("speaker_script_zh") or ""
        if text:
            notes[slide["slide_id"]] = text
    return notes


def add_speaker_notes(slide, text: str) -> None:
    if not text:
        return
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    text_frame.text = text


def items_for_slide(slide_id: str, ocr_dir: Path, text_items_dir: Path | None) -> list[dict]:
    if text_items_dir:
        text_items = text_items_dir / f"{slide_id}_all_ocr_text_items.json"
        if text_items.exists():
            return json.loads(text_items.read_text(encoding="utf-8"))
    return json.loads((ocr_dir / f"{slide_id}_ocr.json").read_text(encoding="utf-8"))


def image_for_slide(slide_id: str, image_dir: Path) -> Path:
    matches = sorted(image_dir.glob(f"{slide_id}*.png")) + sorted(image_dir.glob(f"{slide_id}*.jpg")) + sorted(image_dir.glob(f"{slide_id}*.jpeg"))
    if not matches:
        raise FileNotFoundError(f"No image found for {slide_id} in {image_dir}")
    return matches[0]


def add_full_slide_image(slide, image_path: Path, prs: Presentation) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def build_image_only_pptx(image_dir: Path, out: Path, manifest_path: Path | None, notes_path: Path | None) -> None:
    notes = load_notes(notes_path if notes_path else manifest_path)
    slide_ids = slide_ids_from_manifest(manifest_path)
    if not slide_ids:
        slide_ids = [p.stem for p in sorted(image_dir.glob("slide*.png"))]

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    for slide_id in slide_ids:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_slide_image(slide, image_for_slide(slide_id, image_dir), prs)
        add_speaker_notes(slide, notes.get(slide_id, ""))

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build editable or image-only PPTX files from generated slide pipeline outputs.")
    parser.add_argument("--mode", choices=["editable", "image-only"], default="editable")
    parser.add_argument("--ocr-dir")
    parser.add_argument("--clean-dir")
    parser.add_argument("--image-dir")
    parser.add_argument("--out", required=True)
    parser.add_argument("--manifest")
    parser.add_argument("--text-items-dir")
    parser.add_argument("--notes-json")
    args = parser.parse_args()

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    manifest_path = Path(args.manifest) if args.manifest else None
    notes_path = Path(args.notes_json) if args.notes_json else None

    if args.mode == "image-only":
        if not args.image_dir:
            raise ValueError("--image-dir is required in image-only mode")
        build_image_only_pptx(Path(args.image_dir), out, manifest_path, notes_path)
        print(out)
        return

    if not args.ocr_dir or not args.clean_dir:
        raise ValueError("--ocr-dir and --clean-dir are required in editable mode")

    ocr_dir = Path(args.ocr_dir)
    clean_dir = Path(args.clean_dir)
    text_items_dir = Path(args.text_items_dir) if args.text_items_dir else None
    notes = load_notes(notes_path if notes_path else manifest_path)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    for slide_id in slide_ids_from_manifest(manifest_path, ocr_dir):
        ocr_json = ocr_dir / f"{slide_id}_ocr.json"
        clean = clean_dir / f"{slide_id}_clean_background.png"
        if not ocr_json.exists():
            raise FileNotFoundError(ocr_json)
        if not clean.exists():
            raise FileNotFoundError(clean)
        items = merge_adjacent_heading_items(items_for_slide(slide_id, ocr_dir, text_items_dir))
        source_image = None
        if args.image_dir:
            source_image = Image.open(image_for_slide(slide_id, Path(args.image_dir))).convert("RGB")
        clean_image = Image.open(clean).convert("RGB")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(clean), 0, 0, width=prs.slide_width, height=prs.slide_height)
        for item in items:
            add_textbox(slide, item, source_image, clean_image)
        add_speaker_notes(slide, notes.get(slide_id, ""))

    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
